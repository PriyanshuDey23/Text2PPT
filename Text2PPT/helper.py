import streamlit as st
import base64
import google.generativeai as genai
import pptx
from pptx.util import Pt
import os
from dotenv import load_dotenv

# Load the API key
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    st.error("API key not found. Please set the GOOGLE_API_KEY environment variable.")
    st.stop()

# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)
MAX_CONTENT_LENGTH = 1000  # Max characters per slide content

# Configure the API key for Google Generative AI
try:
    genai.configure(api_key=GOOGLE_API_KEY)
except Exception as e:
    st.error(f"Failed to configure Generative AI: {e}")
    st.stop()

# Function to initialize the generative model
def get_generative_model():
    try:
        # Modify this based on your actual generative model initialization.
        model = genai.GenerativeModel(model_name="gemini-1.5-flash")
        return model
    except Exception as e:
        st.error(f"Error initializing the generative model: {e}")
        st.stop()

model = get_generative_model()

# Function to generate slide titles
def generate_slide_titles(topic, tone, language, number_of_slides):
    try:
        chat_session = model.start_chat()  # This depends on your API structure
        input_message = f"Generate {number_of_slides} slide titles for a presentation on '{topic}' with a {tone} tone in {language}."
        response = chat_session.send_message(input_message)
        if not response or not response.text.strip():
            raise ValueError("Empty response from the generative model.")
        return [title.strip() for title in response.text.split("\n") if title.strip()]
    except Exception as e:
        st.error(f"Error generating slide titles: {e}")
        return []

# Function to generate slide content
def generate_slide_content(slide_title, tone, language):
    try:
        chat_session = model.start_chat()  # This depends on your API structure
        input_message = f"Generate content for the slide: '{slide_title}' in {language} with a {tone} tone."
        response = chat_session.send_message(input_message)
        if not response or not response.text.strip():
            raise ValueError("Empty response from the generative model.")
        return response.text.strip()
    except Exception as e:
        st.error(f"Error generating slide content for '{slide_title}': {e}")
        return ""

# Function to split content into smaller chunks
def split_content(content):
    # Split the content into chunks if it exceeds the max length
    chunks = []
    while len(content) > MAX_CONTENT_LENGTH:
        split_idx = content.rfind(" ", 0, MAX_CONTENT_LENGTH)
        chunks.append(content[:split_idx])
        content = content[split_idx:].strip()
    chunks.append(content)  # Add the remaining content
    return chunks

# Function to create the PowerPoint presentation
def create_presentation(topic, slide_titles, slide_contents):
    try:
        prs = pptx.Presentation()
        slide_layout = prs.slide_layouts[1]

        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = topic

        for slide_title, slide_content in zip(slide_titles, slide_contents):
            content_chunks = split_content(slide_content)  # Split content if needed
            for chunk in content_chunks:
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = slide_title
                slide.shapes.placeholders[1].text = chunk

                # Customize font size for titles and content
                slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            paragraph.font.size = SLIDE_FONT_SIZE

        output_dir = "generated_ppt"
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"{topic}_presentation.pptx")
        prs.save(output_path)
        return output_path
    except Exception as e:
        st.error(f"Error creating presentation: {e}")
        return ""

# Function to provide a download link for the generated presentation
def get_ppt_download_link(file_path):
    try:
        with open(file_path, "rb") as file:
            ppt_contents = file.read()

        b64_ppt = base64.b64encode(ppt_contents).decode()
        return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{os.path.basename(file_path)}">Download the PowerPoint Presentation</a>'
    except Exception as e:
        st.error(f"Error generating download link: {e}")
        return ""